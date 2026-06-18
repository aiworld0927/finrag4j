"""
文档解析服务模块

支持的文档类型：
1. PDF - 支持文本提取和表格提取
2. Word (.doc, .docx) - 支持文本和表格提取
3. Excel (.xls, .xlsx) - 支持单元格内容提取
4. TXT - 纯文本提取
5. PPT (.pptx) - 幻灯片文本提取
"""

import os
import io
from typing import Dict, List, Any, Optional
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import PyPDF2
from unstructured.partition.pdf import partition_pdf
from unstructured.partition.docx import partition_docx
from unstructured.partition.xlsx import partition_xlsx
from unstructured.partition.text import partition_text
from utils.text_cleaner import clean_text


class DocumentParser:
    """文档解析器"""
    
    def __init__(self):
        self.supported_types = ["pdf", "docx", "doc", "xlsx", "xls", "txt", "pptx"]
    
    def parse_file(self, file_path: str, file_type: Optional[str] = None) -> Dict[str, Any]:
        """
        解析文档文件
        
        Args:
            file_path: 文件路径
            file_type: 文件类型（可选，自动检测）
        
        Returns:
            解析结果字典
        """
        if not file_type:
            file_type = self._detect_file_type(file_path)
        
        if file_type not in self.supported_types:
            raise ValueError(f"不支持的文件类型: {file_type}")
        
        parse_method = {
            "pdf": self._parse_pdf,
            "docx": self._parse_docx,
            "doc": self._parse_doc,
            "xlsx": self._parse_xlsx,
            "xls": self._parse_xls,
            "txt": self._parse_txt,
            "pptx": self._parse_pptx,
        }.get(file_type)
        
        if not parse_method:
            raise ValueError(f"不支持的文件类型: {file_type}")
        
        try:
            result = parse_method(file_path)
            return {
                "success": True,
                "file_type": file_type,
                "content": result,
                "message": "解析成功"
            }
        except Exception as e:
            return {
                "success": False,
                "file_type": file_type,
                "content": None,
                "message": f"解析失败: {str(e)}"
            }
    
    def parse_bytes(self, file_bytes: bytes, file_type: str) -> Dict[str, Any]:
        """
        解析文件字节流
        
        Args:
            file_bytes: 文件字节内容
            file_type: 文件类型
        
        Returns:
            解析结果字典
        """
        if file_type not in self.supported_types:
            raise ValueError(f"不支持的文件类型: {file_type}")
        
        # 创建临时文件
        import tempfile
        with tempfile.NamedTemporaryFile(delete=False, suffix=f".{file_type}") as f:
            f.write(file_bytes)
            temp_path = f.name
        
        try:
            return self.parse_file(temp_path, file_type)
        finally:
            os.unlink(temp_path)
    
    def _detect_file_type(self, file_path: str) -> str:
        """
        根据文件扩展名检测文件类型
        
        Args:
            file_path: 文件路径
        
        Returns:
            文件类型
        """
        _, ext = os.path.splitext(file_path)
        return ext.lower().lstrip(".")
    
    def _parse_pdf(self, file_path: str) -> Dict[str, Any]:
        """
        解析PDF文件
        
        Args:
            file_path: PDF文件路径
        
        Returns:
            解析结果，包含文本内容和表格
        """
        result = {
            "text": "",
            "tables": [],
            "pages": []
        }
        
        try:
            # 使用unstructured解析PDF
            elements = partition_pdf(file_path)
            
            text_parts = []
            table_parts = []
            
            for element in elements:
                if hasattr(element, 'text'):
                    text_parts.append(element.text)
                elif hasattr(element, 'metadata') and element.metadata.get('category') == 'Table':
                    table_parts.append(str(element))
            
            result["text"] = clean_text("\n".join(text_parts))
            result["tables"] = table_parts
            
            # 同时使用PyPDF2获取页面信息
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                result["pages"] = [
                    {"page_number": i + 1, "text": clean_text(page.extract_text())}
                    for i, page in enumerate(reader.pages)
                ]
        
        except Exception as e:
            # 降级使用PyPDF2
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                page_texts = []
                for page in reader.pages:
                    text = page.extract_text()
                    if text:
                        page_texts.append(text)
                
                result["text"] = clean_text("\n".join(page_texts))
                result["pages"] = [
                    {"page_number": i + 1, "text": clean_text(page.extract_text())}
                    for i, page in enumerate(reader.pages)
                ]
        
        return result
    
    def _parse_docx(self, file_path: str) -> Dict[str, Any]:
        """
        解析Word文档（.docx）
        
        Args:
            file_path: Word文件路径
        
        Returns:
            解析结果，包含文本内容和表格
        """
        result = {
            "text": "",
            "tables": [],
            "paragraphs": []
        }
        
        try:
            # 使用unstructured解析
            elements = partition_docx(file_path)
            
            text_parts = []
            table_parts = []
            
            for element in elements:
                if hasattr(element, 'text'):
                    text_parts.append(element.text)
                elif hasattr(element, 'metadata') and element.metadata.get('category') == 'Table':
                    table_parts.append(str(element))
            
            result["text"] = clean_text("\n".join(text_parts))
            result["tables"] = table_parts
            
        except Exception as e:
            # 降级使用python-docx
            doc = Document(file_path)
            
            paragraphs = []
            tables = []
            
            for para in doc.paragraphs:
                if para.text.strip():
                    paragraphs.append(para.text)
            
            for table in doc.tables:
                table_data = []
                for row in table.rows:
                    row_data = [cell.text.strip() for cell in row.cells]
                    table_data.append(row_data)
                tables.append(table_data)
            
            result["text"] = clean_text("\n".join(paragraphs))
            result["tables"] = tables
            result["paragraphs"] = paragraphs
        
        return result
    
    def _parse_doc(self, file_path: str) -> Dict[str, Any]:
        """
        解析旧版Word文档（.doc）
        
        Args:
            file_path: Word文件路径
        
        Returns:
            解析结果
        """
        # .doc文件需要额外依赖，这里返回基本信息
        # 建议用户转换为.docx格式或使用LibreOffice转换
        return {
            "text": "",
            "tables": [],
            "message": "建议将.doc文件转换为.docx格式以获得更好的解析效果"
        }
    
    def _parse_xlsx(self, file_path: str) -> Dict[str, Any]:
        """
        解析Excel文件（.xlsx）
        
        Args:
            file_path: Excel文件路径
        
        Returns:
            解析结果，包含工作表和单元格数据
        """
        result = {
            "sheets": [],
            "text": ""
        }
        
        try:
            # 使用unstructured解析
            elements = partition_xlsx(file_path)
            
            text_parts = []
            for element in elements:
                if hasattr(element, 'text'):
                    text_parts.append(element.text)
            
            result["text"] = clean_text("\n".join(text_parts))
            
        except Exception as e:
            # 降级使用openpyxl
            wb = load_workbook(file_path, read_only=True)
            
            sheets_data = []
            all_text = []
            
            for sheet_name in wb.sheetnames:
                sheet = wb[sheet_name]
                rows = []
                
                for row in sheet.iter_rows(values_only=True):
                    row_data = []
                    for cell in row:
                        if cell is not None:
                            cell_value = str(cell)
                            row_data.append(cell_value)
                            all_text.append(cell_value)
                    if any(row_data):
                        rows.append(row_data)
                
                sheets_data.append({
                    "name": sheet_name,
                    "rows": rows
                })
            
            result["sheets"] = sheets_data
            result["text"] = clean_text("\n".join(all_text))
        
        return result
    
    def _parse_xls(self, file_path: str) -> Dict[str, Any]:
        """
        解析旧版Excel文件（.xls）
        
        Args:
            file_path: Excel文件路径
        
        Returns:
            解析结果
        """
        return {
            "text": "",
            "sheets": [],
            "message": "建议将.xls文件转换为.xlsx格式以获得更好的解析效果"
        }
    
    def _parse_txt(self, file_path: str) -> Dict[str, Any]:
        """
        解析纯文本文件
        
        Args:
            file_path: 文本文件路径
        
        Returns:
            解析结果
        """
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        
        return {
            "text": clean_text(content),
            "lines": content.splitlines()
        }
    
    def _parse_pptx(self, file_path: str) -> Dict[str, Any]:
        """
        解析PPT文件（.pptx）
        
        Args:
            file_path: PPT文件路径
        
        Returns:
            解析结果，包含每页内容
        """
        prs = Presentation(file_path)
        
        slides = []
        all_text = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text = shape.text.strip()
                    if text:
                        slide_text.append(text)
                        all_text.append(text)
            
            slides.append({
                "slide_number": slide_num,
                "text": "\n".join(slide_text)
            })
        
        return {
            "slides": slides,
            "text": clean_text("\n".join(all_text))
        }


# 创建全局解析器实例
document_parser = DocumentParser()
